import base64
import ctypes
import datetime
import glob
import json
import os
import random
import re
import shutil
import sys  # <--- 这一行绝对不能少
import time
import zipfile

import winreg
from PyQt6.QtCore import (
    QUrl, QObject, pyqtSlot, pyqtSignal, QTimer, Qt, QSharedMemory,QRect,QSize,QPoint
)
from PyQt6.QtGui import QColor, QIcon, QAction, QCursor, QDesktopServices, QPixmap, QPainter, QPen, QTransform, \
    QKeySequence,QShortcut
from PyQt6.QtWebChannel import QWebChannel
from PyQt6.QtWebEngineCore import (
    QWebEngineProfile, QWebEnginePage, QWebEngineSettings
)
from PyQt6.QtWebEngineWidgets import QWebEngineView
from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget,
    QSystemTrayIcon, QMenu, QDialog, QGraphicsView, QGraphicsScene, QGraphicsPixmapItem,
    QVBoxLayout, QHBoxLayout, QPushButton, QColorDialog, QFileDialog,
    QRubberBand
)

from config import (
    NOTES_DIR, IMAGES_DIR, ATTACHMENTS_DIR, TRASH_DIR, STAGING_DIR,
    ORDER_FILE, TRASH_FILE, MAX_TRASH_SIZE
)
# ---------- 导入自定义模块 ----------
from utils import (
    get_abs_path, file_url_from_path, safe_remove,
    read_json_file, write_json_file, is_windows
)

# ---- 暂不导入 keyboard / mss / PIL，防止崩溃 ----
# import keyboard
# import mss
# from PIL import Image

import ctypes
from ctypes import wintypes

def _get_windows_font_families():
    """返回 Windows 下所有字体的本地化族名（中文系统即中文名）"""
    LF_FACESIZE = 32
    class LOGFONTW(ctypes.Structure):
        _fields_ = [
            ("lfHeight", wintypes.LONG),
            ("lfWidth", wintypes.LONG),
            ("lfEscapement", wintypes.LONG),
            ("lfOrientation", wintypes.LONG),
            ("lfWeight", wintypes.LONG),
            ("lfItalic", wintypes.BYTE),
            ("lfUnderline", wintypes.BYTE),
            ("lfStrikeOut", wintypes.BYTE),
            ("lfCharSet", wintypes.BYTE),
            ("lfOutPrecision", wintypes.BYTE),
            ("lfClipPrecision", wintypes.BYTE),
            ("lfQuality", wintypes.BYTE),
            ("lfPitchAndFamily", wintypes.BYTE),
            ("lfFaceName", wintypes.WCHAR * LF_FACESIZE),
        ]

    class TEXTMETRICW(ctypes.Structure):
        _fields_ = [
            ("tmHeight", wintypes.LONG),
            ("tmAscent", wintypes.LONG),
            ("tmDescent", wintypes.LONG),
            ("tmInternalLeading", wintypes.LONG),
            ("tmExternalLeading", wintypes.LONG),
            ("tmAveCharWidth", wintypes.LONG),
            ("tmMaxCharWidth", wintypes.LONG),
            ("tmWeight", wintypes.LONG),
            ("tmOverhang", wintypes.LONG),
            ("tmDigitizedAspectX", wintypes.LONG),
            ("tmDigitizedAspectY", wintypes.LONG),
            ("tmFirstChar", wintypes.WCHAR),
            ("tmLastChar", wintypes.WCHAR),
            ("tmDefaultChar", wintypes.WCHAR),
            ("tmBreakChar", wintypes.WCHAR),
            ("tmItalic", wintypes.BYTE),
            ("tmUnderlined", wintypes.BYTE),
            ("tmStruckOut", wintypes.BYTE),
            ("tmPitchAndFamily", wintypes.BYTE),
            ("tmCharSet", wintypes.BYTE),
        ]

    FONTENUMPROC = ctypes.WINFUNCTYPE(
        wintypes.INT,
        ctypes.POINTER(LOGFONTW),
        ctypes.POINTER(TEXTMETRICW),
        wintypes.DWORD,
        wintypes.LPARAM
    )

    fonts = set()
    def callback(lpelf, lpntm, font_type, lparam):
        name = str(lpelf.contents.lfFaceName)
        if name:
            clean = name.rstrip('\x00')
            if not clean.startswith('@'):
                fonts.add(clean)
        return 1

    hdc = ctypes.windll.user32.GetDC(0)
    lf = LOGFONTW()
    lf.lfCharSet = 1          # DEFAULT_CHARSET
    lf.lfFaceName = '\0' * LF_FACESIZE
    lf.lfPitchAndFamily = 0

    ctypes.windll.gdi32.EnumFontFamiliesExW(
        hdc,
        ctypes.byref(lf),
        FONTENUMPROC(callback),
        0,                      # lparam 不需要，直接用闭包
        0
    )
    ctypes.windll.user32.ReleaseDC(0, hdc)
    return fonts

# ---------- 危险环境变量已移除 ----------
# os.environ['QT_OPENGL'] = 'angle'
# os.environ['QTWEBENGINE_CHROMIUM_FLAGS'] = '--ignore-gpu-blocklist --disable-gpu-driver-bug-workarounds'
# os.environ['QTWEBENGINE_REMOTE_DEBUGGING'] = '9222'

# ---------- 全局变量 ----------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ---------- 图片编辑器 ----------
class ImageEditorDialog(QDialog):
    def __init__(self, image_path, parent=None):
        super().__init__(parent)
        self.setWindowTitle("图片编辑器")
        self.resize(900, 700)
        self.image_path = image_path
        self.toast_message_signal = None

        if not os.path.exists(self.image_path):
            self.original_pixmap = QPixmap()
            if self.toast_message_signal:
                self.toast_message_signal.emit("⚠️ 找不到该图片文件，请检查路径！")
        else:
            self.original_pixmap = QPixmap(self.image_path)
            if self.original_pixmap.isNull():
                if self.toast_message_signal:
                    self.toast_message_signal.emit("⚠️ 图片格式损坏或无法解析！")

        self.edited_pixmap = self.original_pixmap.copy()
        self.rotation_angle = 0
        self.current_tool = "pen"
        self.is_drawing = False
        self.last_point = None

        self.scene = QGraphicsScene()
        self.pixmap_item = QGraphicsPixmapItem(self.edited_pixmap)
        self.scene.addItem(self.pixmap_item)
        self.view = QGraphicsView(self.scene)
        self.view.setStyleSheet("background: #f0f0f0;")

        toolbar = QHBoxLayout()
        self.btn_pen = QPushButton("🖊 画笔")
        self.btn_eraser = QPushButton("🧹 橡皮擦")
        self.btn_color = QPushButton("🎨 选颜色")
        self.btn_rotate = QPushButton("🔄 旋转 90°")
        self.btn_save = QPushButton("💾 保存")
        self.btn_cancel = QPushButton("❌ 取消")

        self.btn_pen.clicked.connect(lambda: self.set_tool("pen"))
        self.btn_eraser.clicked.connect(lambda: self.set_tool("eraser"))
        self.btn_color.clicked.connect(self.choose_color)
        self.btn_rotate.clicked.connect(self.rotate_image)
        self.btn_save.clicked.connect(self.save_and_close)
        self.btn_cancel.clicked.connect(self.reject)

        toolbar.addWidget(self.btn_pen)
        toolbar.addWidget(self.btn_eraser)
        toolbar.addWidget(self.btn_color)
        toolbar.addWidget(self.btn_rotate)
        toolbar.addStretch()
        toolbar.addWidget(self.btn_save)
        toolbar.addWidget(self.btn_cancel)

        layout = QVBoxLayout()
        layout.addLayout(toolbar)
        layout.addWidget(self.view)
        self.setLayout(layout)

        self.color = QColor(0, 0, 0)
        self.eraser_size = 20

    def set_tool(self, tool):
        self.current_tool = tool
        self.view.setDragMode(QGraphicsView.DragMode.NoDrag)

    def choose_color(self):
        color = QColorDialog.getColor()
        if color.isValid():
            self.color = color

    def rotate_image(self):
        self.rotation_angle = (self.rotation_angle + 90) % 360
        transform = QTransform().rotate(self.rotation_angle)
        self.edited_pixmap = self.edited_pixmap.transformed(transform)
        self.pixmap_item.setPixmap(self.edited_pixmap)
        self.view.setSceneRect(0, 0, self.edited_pixmap.width(), self.edited_pixmap.height())

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            self.is_drawing = True
            self.last_point = self.view.mapToScene(event.pos()).toPoint()

    def mouseMoveEvent(self, event):
        if self.is_drawing and self.last_point:
            current_point = self.view.mapToScene(event.pos()).toPoint()
            painter = QPainter(self.edited_pixmap)
            if self.current_tool == "pen":
                painter.setPen(QPen(self.color, 3, Qt.PenStyle.SolidLine, Qt.PenCapStyle.RoundCap, Qt.PenJoinStyle.RoundJoin))
            elif self.current_tool == "eraser":
                painter.setCompositionMode(QPainter.CompositionMode.CompositionMode_Clear)
                painter.setPen(QPen(Qt.GlobalColor.transparent, self.eraser_size,
                                    Qt.PenStyle.SolidLine, Qt.PenCapStyle.RoundCap, Qt.PenJoinStyle.RoundJoin))
            painter.drawLine(self.last_point, current_point)
            painter.end()
            self.pixmap_item.setPixmap(self.edited_pixmap)
            self.last_point = current_point

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            self.is_drawing = False
            self.last_point = None

    def wheelEvent(self, event):
        zoomInFactor = 1.1
        zoomOutFactor = 1 / zoomInFactor
        if event.angleDelta().y() > 0:
            self.view.scale(zoomInFactor, zoomInFactor)
        else:
            self.view.scale(zoomOutFactor, zoomOutFactor)

    def save_and_close(self):
        save_path, _ = QFileDialog.getSaveFileName(self, "保存图片", os.path.dirname(self.image_path),
                                                   "Images (*.png *.jpg)")
        if save_path:
            self.edited_pixmap.save(save_path)
            if self.toast_message_signal:
                self.toast_message_signal.emit(f"✅ 编辑后的图片已保存至: {save_path}")
            self.accept()

class ScreenCaptureWidget(QWidget):
    """全屏选区工具：仅绘制遮罩与选框，完成选区后通过信号传递矩形并关闭"""
    region_selected = pyqtSignal(QRect)   # 发送选区矩形（全局坐标）

    def __init__(self):
        super().__init__()
        self.setWindowFlags(
            Qt.WindowType.FramelessWindowHint |
            Qt.WindowType.WindowStaysOnTopHint |
            Qt.WindowType.Tool
        )
        self.setAttribute(Qt.WidgetAttribute.WA_TranslucentBackground, True)
        # 通过 paintEvent 绘制半透明背景，避免样式表与 WA_TranslucentBackground 冲突
        self.setAutoFillBackground(False)

        screen = QApplication.primaryScreen()
        if screen:
            self.setGeometry(screen.geometry())

        self.origin = QPoint()
        self.current_rect = QRect()

    def paintEvent(self, event):
        painter = QPainter(self)
        # 绘制半透明黑色遮罩
        painter.fillRect(self.rect(), QColor(0, 0, 0, 80))
        # 如果有有效选区，绘制选框（白色边框）
        if self.current_rect.isValid() and self.current_rect.width() > 0 and self.current_rect.height() > 0:
            painter.setPen(QPen(QColor(255, 255, 255), 2))
            painter.drawRect(self.current_rect)

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            self.origin = event.pos()
            self.current_rect = QRect(self.origin, QSize())
            self.update()

    def mouseMoveEvent(self, event):
        if event.buttons() & Qt.MouseButton.LeftButton:
            self.current_rect = QRect(self.origin, event.pos()).normalized()
            self.update()

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            rect = self.current_rect
            if rect.width() > 10 and rect.height() > 10:
                # 发射全局坐标矩形
                self.region_selected.emit(rect)
            self.close()

    def keyPressEvent(self, event):
        if event.key() == Qt.Key.Key_Escape:
            self.close()

# ---------- 自定义 WebEnginePage ----------
class CustomWebEnginePage(QWebEnginePage):
    def acceptNavigationRequest(self, url, navigation_type, is_main_frame):
        if url.scheme() in ["http", "https"]:
            QDesktopServices.openUrl(url)
            return False
        if url.scheme() == "file":
            path = url.path()
            if "attachments" in path or path.endswith(('.txt', '.md', '.log', '.png', '.jpg', '.pdf', '.docx')):
                QDesktopServices.openUrl(url)
                return False
        return super().acceptNavigationRequest(url, navigation_type, is_main_frame)

# ---------- FloatActions 单独给浮窗调用的动作 ----------
class FloatActions(QObject):
    toggle_top_signal = pyqtSignal()
    open_editor_signal = pyqtSignal(str)

    @pyqtSlot(str, str, result=str)
    def save_file(self, file_name, base64_data):
        return self.floating_window.backend.save_float_file(file_name, base64_data)

    def __init__(self, floating_window):
        super().__init__()
        self.floating_window = floating_window
        self.toggle_top_signal.connect(floating_window.toggle_float_top)
        self.open_editor_signal.connect(floating_window.open_image_editor)

    @pyqtSlot()
    def toggle_top(self):
        self.toggle_top_signal.emit()

    @pyqtSlot(str)
    def open_image_editor(self, image_path):
        self.open_editor_signal.emit(image_path)

# ---------- Backend 核心后端 ----------
class Backend(QObject):
    image_saved_signal = pyqtSignal(str)
    update_sidebar_signal = pyqtSignal(str)
    load_note_signal = pyqtSignal(str, str)
    update_trash_signal = pyqtSignal(str)
    load_fonts_signal = pyqtSignal(str)
    theme_packs_signal = pyqtSignal(str)
    toast_message_signal = pyqtSignal(str)
    attachment_saved_signal = pyqtSignal(str)
    text_imported_signal = pyqtSignal(str)
    native_menu_action = pyqtSignal(str, str)

    def __init__(self, parent_window):
        super().__init__()
        self.parent_window = parent_window
        self.detached_note_ids = set()   # 支持多个浮窗
        self._sidebar_refresh_timer = QTimer()
        self._sidebar_refresh_timer.setSingleShot(True)
        self._sidebar_refresh_timer.timeout.connect(self.get_all_notes)

    @pyqtSlot(str, bool)
    def open_native_context_menu(self, note_id, is_pinned):
        menu = QMenu(self.parent_window)
        actions = [
            ("new_note", "📝 新建笔记"),
            ("duplicate", "📋 复制笔记"),
            ("float", "💨 分离到磁贴"),
            ("pin", "📌 取消置顶" if is_pinned else "📌 置顶当前笔记"),
            ("export_md", "📤 导出笔记 (另存为MD)"),
            ("export_pdf", "📄 导出为 PDF"),
            ("delete", "🗑️ 删除当前笔记"),
        ]
        for action_id, text in actions:
            act = menu.addAction(text)
            act.triggered.connect(lambda checked, aid=action_id: self.native_menu_action.emit(note_id, aid))
        menu.exec(QCursor.pos())

    @pyqtSlot(str)
    def paste_image_from_clipboard(self, base64_data):
        try:
            header, encoded = base64_data.split(',', 1)
            ext = "png"
            if "jpeg" in header or "jpg" in header:
                ext = "jpg"
            elif "png" in header:
                ext = "png"
            file_name = f"pasted_{int(time.time())}.{ext}"
            save_path = get_abs_path(IMAGES_DIR, file_name)
            with open(save_path, "wb") as f:
                f.write(base64.b64decode(encoded))
            self.image_saved_signal.emit(save_path.replace("\\", "/"))
            self.toast_message_signal.emit("图片粘贴保存成功！")
        except Exception as e:
            print(f"粘贴图片失败: {e}")
            self.toast_message_signal.emit("粘贴失败！")

    @pyqtSlot(str, str)
    def save_attached_file(self, file_name, base64_data):
        try:
            header, encoded = base64_data.split(',', 1)
            save_path = get_abs_path(ATTACHMENTS_DIR, file_name)
            with open(save_path, "wb") as f:
                f.write(base64.b64decode(encoded))

            # 中转站备份
            try:
                staging_path = get_abs_path(STAGING_DIR, file_name)
                if not os.path.exists(staging_path):
                    shutil.copy2(save_path, staging_path)
            except Exception as e:
                print(f"⚠️ 中转站备份附件失败: {e}")

            lower_name = file_name.lower()

            if lower_name.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                img_save_path = get_abs_path(IMAGES_DIR, file_name)
                with open(img_save_path, "wb") as f:
                    f.write(base64.b64decode(encoded))

                # 【绝杀修复】使用 urllib.parse.quote 强制将整条路径所有特殊字符都转义
                import urllib.parse
                # 1. 将反斜杠转正斜杠
                safe_path = img_save_path.replace("\\", "/")
                # 2. 全路径 URL 编码 (空格变%20, 括号变%28%29, 中文也会自动编码)
                full_url = "file:///" + urllib.parse.quote(safe_path)

                self.image_saved_signal.emit(full_url)
                self.toast_message_signal.emit(f"图片已保存: {file_name}")
                return

            txt_extensions = ('.txt', '.md', '.log', '.csv', '.xml', '.json', '.yaml', '.yml', '.ini', '.conf',
                              '.py', '.cmd', '.bat', '.sh', '.ps1', '.vbs', '.js', '.lua', '.r', '.pl', '.java', '.c',
                              '.cpp', '.h')
            if lower_name.endswith(txt_extensions):
                try:
                    with open(save_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    with open(save_path, 'r', encoding='gbk', errors='ignore') as f:
                        content = f.read()
                if len(content) > 50000:
                    content = content[:50000] + "\n... (文本过长已截断)"
                ext = os.path.splitext(file_name)[1][1:]
                self.text_imported_signal.emit(f"\n```{ext}\n{content}\n```\n")
                self.toast_message_signal.emit(f"已导入代码/配置文件: {file_name}")
                return

            if lower_name.endswith('.pdf'):
                try:
                    import fitz
                    doc = fitz.open(save_path)
                    text_content = ""
                    for page_num, page in enumerate(doc):
                        text_content += page.get_text("text") + "\n\n"
                        for img_index, img in enumerate(page.get_images(full=True)):
                            try:
                                xref = img[0]
                                pix = fitz.Pixmap(doc, xref)
                                if pix.n - pix.alpha < 4:
                                    pix = fitz.Pixmap(fitz.csRGB, pix)
                                img_filename = f"pdf_img_{int(time.time())}_{page_num}_{img_index}.png"
                                img_path = get_abs_path(IMAGES_DIR, img_filename)
                                pix.save(img_path)
                                # 同理，PDF 的图片路径也做完整编码
                                import urllib.parse
                                safe_path = img_path.replace("\\", "/")
                                full_url = "file:///" + urllib.parse.quote(safe_path)
                                self.image_saved_signal.emit(full_url)
                            except Exception as e:
                                print(f"PDF 某张图片提取失败: {e}")
                    doc.close()
                    if len(text_content) > 50000:
                        text_content = text_content[:50000] + "\n... (PDF文本过长已截断)"
                    if not text_content.strip():
                        self.toast_message_signal.emit("⚠️ 该 PDF 是纯图片，无文字。")
                    else:
                        self.text_imported_signal.emit(f"\n```text\n[PDF预览]\n{text_content}\n```\n")
                        self.toast_message_signal.emit(f"已导入 {file_name} 的 PDF 内容及图片！")
                        return
                except Exception as e:
                    print(f"❌【PDF解析错误】: {e}")
                    self.toast_message_signal.emit("❌ PDF 解析失败，将在外部打开。")

            if lower_name.endswith('.docx'):
                try:
                    from docx import Document
                    doc = Document(save_path)
                    text = "\n".join([para.text for para in doc.paragraphs])
                    if len(text) > 50000:
                        text = text[:50000] + "\n... (Word文本过长已截断)"
                    self.text_imported_signal.emit(f"\n```text\n[Word预览]\n{text}\n```\n")
                    self.toast_message_signal.emit(f"已导入 {file_name} 的 Word 文字预览！")
                    return
                except Exception as e:
                    print(f"Word解析失败: {e}")
                    self.toast_message_signal.emit("❌ Word 解析失败，将在外部打开。")

            if lower_name.endswith('.pptx'):
                try:
                    from pptx import Presentation
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    prs = Presentation(save_path)
                    text_content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                table = shape.table
                                for row in table.rows:
                                    row_cells = [cell.text.strip() for cell in row.cells]
                                    text_content += " | ".join(row_cells) + "\n"
                            elif hasattr(shape, "text") and shape.text:
                                text_content += shape.text + "\n"
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                try:
                                    image = shape.image
                                    img_filename = f"pptx_img_{int(time.time())}_{random.randint(100, 999)}.png"
                                    img_path = get_abs_path(IMAGES_DIR, img_filename)
                                    with open(img_path, "wb") as f:
                                        f.write(image.blob)
                                    self.image_saved_signal.emit(img_path.replace("\\", "/"))
                                except Exception as e:
                                    print(f"PPTX 图片提取失败: {e}")
                    if not text_content.strip():
                        self.toast_message_signal.emit("⚠️ 该 PPTX 中没有可提取的文本。")
                    else:
                        if len(text_content) > 50000:
                            text_content = text_content[:50000] + "\n... (PPT文本过长已截断)"
                        self.text_imported_signal.emit(f"\n```text\n[PPTX预览]\n{text_content}\n```\n")
                        self.toast_message_signal.emit(f"已导入 {file_name} 的 PPTX 全部文本、表格和图片！")
                        return
                except Exception as e:
                    print(f"PPTX解析失败: {e}")
                    self.toast_message_signal.emit("❌ PPTX 解析失败，将在外部打开。")

            if lower_name.endswith(('.xlsx', '.xls')):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(save_path, data_only=True)
                    text_content = f"📊 工作表: {', '.join(wb.sheetnames)}\n\n"
                    for sheet_index, sheet in enumerate(wb.worksheets):
                        text_content += f"--- {sheet.title} ---\n"
                        rows = list(sheet.iter_rows(values_only=True))
                        if not rows: continue
                        max_cols = max([len(row) for row in rows])
                        headers = [str(rows[0][i]) if i < len(rows[0]) and rows[0][i] is not None else ""
                                   for i in range(max_cols)]
                        text_content += "| " + " | ".join(h.replace('|', '\\|') for h in headers) + " |\n"
                        text_content += "|" + "|".join(["---"] * max_cols) + "|\n"
                        for row in rows[1:31]:
                            row_cells = [str(row[i]).replace('|', '\\|') if i < len(row) and row[i] is not None else ""
                                         for i in range(max_cols)]
                            text_content += "| " + " | ".join(row_cells) + " |\n"
                        if len(rows) > 30:
                            text_content += "... (表格过长，仅显示前30行)\n"
                        text_content += "\n"
                    if len(text_content) > 50000:
                        text_content = text_content[:50000] + "\n...(表格内容过长已截断)"
                    self.text_imported_signal.emit(f"\n\n{text_content}\n\n")
                    self.toast_message_signal.emit(f"已导入 {file_name} 的 Excel 表格预览！")
                    return
                except Exception as e:
                    print(f"Excel解析失败: {e}")
                    self.toast_message_signal.emit("❌ Excel 解析失败，将在外部打开。")

            markdown_link = f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"
            self.attachment_saved_signal.emit(markdown_link)
            self.toast_message_signal.emit(f"附件已添加，点击链接将使用外部软件打开: {file_name}")

        except Exception as e:
            print(f"保存附件失败: {e}")
            self.toast_message_signal.emit("附件保存失败！")

    def save_float_file(self, file_name, base64_data):
        """浮窗专用文件保存，解析文本并提取图片，返回完整 Markdown 字符串"""
        try:
            header, encoded = base64_data.split(',', 1)
            save_path = get_abs_path(ATTACHMENTS_DIR, file_name)
            with open(save_path, "wb") as f:
                f.write(base64.b64decode(encoded))

            # 中转站备份
            staging_path = get_abs_path(STAGING_DIR, file_name)
            if not os.path.exists(staging_path):
                shutil.copy2(save_path, staging_path)

            lower_name = file_name.lower()
            img_links = []  # 收集提取的图片 Markdown

            # ---------- 图片 ----------
            if lower_name.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                img_save_path = get_abs_path(IMAGES_DIR, file_name)
                with open(img_save_path, "wb") as f:
                    f.write(base64.b64decode(encoded))
                url = file_url_from_path(img_save_path)
                return f"![{file_name}]({url})"

            # ---------- 纯文本文件 ----------
            txt_ext = ('.txt', '.md', '.log', '.csv', '.xml', '.json', '.yaml', '.yml',
                       '.ini', '.conf', '.py', '.cmd', '.bat', '.sh', '.ps1', '.vbs',
                       '.js', '.lua', '.r', '.pl', '.java', '.c', '.cpp', '.h')
            if lower_name.endswith(txt_ext):
                try:
                    with open(save_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    with open(save_path, 'r', encoding='gbk', errors='ignore') as f:
                        content = f.read()
                if len(content) > 50000:
                    content = content[:50000] + "\n... (文本过长已截断)"
                ext = os.path.splitext(file_name)[1][1:]
                return f"\n```{ext}\n{content}\n```\n"

            # ---------- PDF（提取文字 + 图片） ----------
            if lower_name.endswith('.pdf'):
                try:
                    import fitz
                    doc = fitz.open(save_path)
                    text_content = ""
                    for page_num, page in enumerate(doc):
                        text_content += page.get_text("text") + "\n\n"
                        # 提取页面图片
                        for img_index, img in enumerate(page.get_images(full=True)):
                            try:
                                xref = img[0]
                                pix = fitz.Pixmap(doc, xref)
                                if pix.n - pix.alpha < 4:
                                    pix = fitz.Pixmap(fitz.csRGB, pix)
                                img_filename = f"pdf_img_{int(time.time())}_{page_num}_{img_index}.png"
                                img_path = get_abs_path(IMAGES_DIR, img_filename)
                                pix.save(img_path)
                                img_links.append(f"![PDF图片]({file_url_from_path(img_path)})")
                            except Exception as e:
                                print(f"PDF图片提取失败: {e}")
                    doc.close()
                    if len(text_content) > 50000:
                        text_content = text_content[:50000] + "\n... (PDF文本过长已截断)"
                    if not text_content.strip() and not img_links:
                        return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"
                    return "\n".join([f"```text\n[PDF预览]\n{text_content}\n```"] + img_links)
                except Exception as e:
                    print(f"PDF解析失败: {e}")
                    return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"

            # ---------- Word（提取文字 + 图片） ----------
            if lower_name.endswith('.docx'):
                try:
                    from docx import Document
                    doc = Document(save_path)
                    text_parts = []
                    for para in doc.paragraphs:
                        text_parts.append(para.text)
                    # 提取内嵌图片
                    for rel in doc.part.rels.values():
                        if "image" in rel.reltype:
                            try:
                                img = rel.target_part.blob
                                ext = rel.target_part.partname.split('.')[-1]
                                if ext.lower() not in ('png', 'jpg', 'jpeg', 'gif', 'bmp'):
                                    ext = 'png'
                                img_filename = f"docx_img_{int(time.time())}_{random.randint(100,999)}.{ext}"
                                img_path = get_abs_path(IMAGES_DIR, img_filename)
                                with open(img_path, 'wb') as f:
                                    f.write(img)
                                img_links.append(f"![Word图片]({file_url_from_path(img_path)})")
                            except Exception as e:
                                print(f"Word图片提取失败: {e}")
                    text = "\n".join(text_parts)
                    if len(text) > 50000:
                        text = text[:50000] + "\n... (Word文本过长已截断)"
                    if not text.strip() and not img_links:
                        return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"
                    return "\n".join([f"```text\n[Word预览]\n{text}\n```"] + img_links)
                except Exception as e:
                    print(f"Word解析失败: {e}")
                    return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"

            # ---------- PPTX（提取文字 + 图片） ----------
            if lower_name.endswith('.pptx'):
                try:
                    from pptx import Presentation
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    prs = Presentation(save_path)
                    text_content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                table = shape.table
                                for row in table.rows:
                                    row_cells = [cell.text.strip() for cell in row.cells]
                                    text_content += " | ".join(row_cells) + "\n"
                            elif hasattr(shape, "text") and shape.text:
                                text_content += shape.text + "\n"
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                try:
                                    image = shape.image
                                    ext = image.content_type.split('/')[-1]
                                    if ext.lower() not in ('png', 'jpg', 'jpeg', 'gif', 'bmp'):
                                        ext = 'png'
                                    img_filename = f"pptx_img_{int(time.time())}_{random.randint(100,999)}.{ext}"
                                    img_path = get_abs_path(IMAGES_DIR, img_filename)
                                    with open(img_path, "wb") as f:
                                        f.write(image.blob)
                                    img_links.append(f"![PPT图片]({file_url_from_path(img_path)})")
                                except Exception as e:
                                    print(f"PPTX图片提取失败: {e}")
                    if len(text_content) > 50000:
                        text_content = text_content[:50000] + "\n... (PPT文本过长已截断)"
                    if not text_content.strip() and not img_links:
                        return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"
                    return "\n".join([f"```text\n[PPTX预览]\n{text_content}\n```"] + img_links)
                except Exception as e:
                    print(f"PPTX解析失败: {e}")
                    return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"

            # ---------- Excel ----------
            if lower_name.endswith(('.xlsx', '.xls')):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(save_path, data_only=True)
                    text_content = f"📊 工作表: {', '.join(wb.sheetnames)}\n\n"
                    for sheet in wb.worksheets:
                        text_content += f"--- {sheet.title} ---\n"
                        rows = list(sheet.iter_rows(values_only=True))
                        if not rows:
                            continue
                        max_cols = max([len(row) for row in rows])
                        headers = [str(rows[0][i]) if i < len(rows[0]) and rows[0][i] is not None else ""
                                   for i in range(max_cols)]
                        text_content += "| " + " | ".join(h.replace('|', '\\|') for h in headers) + " |\n"
                        text_content += "|" + "|".join(["---"] * max_cols) + "|\n"
                        for row in rows[1:31]:
                            row_cells = [str(row[i]).replace('|', '\\|') if i < len(row) and row[i] is not None else ""
                                         for i in range(max_cols)]
                            text_content += "| " + " | ".join(row_cells) + " |\n"
                        if len(rows) > 30:
                            text_content += "... (表格过长，仅显示前30行)\n"
                        text_content += "\n"
                    if len(text_content) > 50000:
                        text_content = text_content[:50000] + "\n...(表格内容过长已截断)"
                    return f"\n\n{text_content}\n\n"
                except Exception as e:
                    print(f"Excel解析失败: {e}")
                    return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"

            # ---------- 其他附件 ----------
            return f"[📎 {file_name}]({save_path.replace(os.sep, '/')})"
        except Exception as e:
            print(f"浮窗文件保存失败: {e}")
            return ""

    @pyqtSlot(str)
    def open_image_editor(self, image_path):
        local_path = QUrl(image_path).toLocalFile()
        abs_path = os.path.abspath(local_path)
        if not os.path.exists(abs_path):
            self.toast_message_signal.emit(f"❌ 找不到该图片文件: {os.path.basename(abs_path)}，可能已被移动或删除！")
            return
        dialog = ImageEditorDialog(abs_path, self.parent_window)
        dialog.toast_message_signal = self.toast_message_signal
        dialog.exec()

    @pyqtSlot()
    def open_staging_folder(self):
        staging_dir = get_abs_path(STAGING_DIR)
        if os.path.exists(staging_dir):
            QDesktopServices.openUrl(QUrl.fromLocalFile(staging_dir))
        else:
            self.toast_message_signal.emit("⚠️ 中转站文件夹不存在！")

    @pyqtSlot(str, str)
    def rename_note(self, note_id, new_title):
        try:
            file_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                if lines:
                    first = lines[0].strip()
                    match = re.match(r'^(#+)', first)
                    if match:
                        lines[0] = f"{match.group(1)} {new_title}\n"
                    else:
                        lines.insert(0, f"# {new_title}\n")
                else:
                    lines = [f"# {new_title}\n"]
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.writelines(lines)
                self.get_all_notes()
                self.get_note_content(note_id)
        except Exception as e:
            print(f"重命名笔记出错: {e}")

    @pyqtSlot(str)
    def export_note_to_file(self, note_id):
        try:
            file_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            if not os.path.exists(file_path):
                self.toast_message_signal.emit("❌ 笔记文件已丢失！")
                return
            with open(file_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
                title = "无标题笔记"
                if lines and lines[0].startswith('# '):
                    title = lines[0][2:].strip()
            save_path, _ = QFileDialog.getSaveFileName(None, "导出笔记", f"{title}.md",
                                                       "Markdown Files (*.md);;All Files (*)")
            if save_path:
                shutil.copy(file_path, save_path)
                self.toast_message_signal.emit(f"✅ 笔记已导出至: {save_path}")
        except Exception as e:
            print(f"导出笔记失败: {e}")

    @pyqtSlot(str)
    def export_note_to_pdf(self, note_id):
        try:
            import pdfkit
            import markdown
            file_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            if not os.path.exists(file_path):
                self.toast_message_signal.emit("❌ 笔记文件已丢失！")
                return
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            lines = content.split('\n')
            title = "无标题笔记"
            if lines and lines[0].startswith('# '):
                title = lines[0][2:].strip()
            html_body = markdown.markdown(content)
            full_html = f"""
            <html><head><meta charset="utf-8"><style>
            body {{ font-family: -apple-system, 'Helvetica Neue', sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #1f1f1f; }}
            table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
            th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
            th {{ background: #f5f5f5; }}
            pre {{ background: #f5f5f5; padding: 10px; border-radius: 4px; overflow-x: auto; }}
            code {{ font-family: 'Courier New', monospace; }}
            </style></head><body>
            {html_body}
            </body></html>
            """
            save_path, _ = QFileDialog.getSaveFileName(None, "导出笔记为 PDF", f"{title}.pdf", "PDF Files (*.pdf)")
            if save_path:
                if getattr(sys, 'frozen', False):
                    base_path = sys._MEIPASS
                else:
                    base_path = os.getcwd()
                wk_path = os.path.join(base_path, 'wkhtml', 'wkhtmltopdf.exe')
                config = pdfkit.configuration(wkhtmltopdf=wk_path)
                pdfkit.from_string(full_html, save_path, configuration=config)
                self.toast_message_signal.emit(f"✅ PDF 已成功导出至: {save_path}")
        except Exception as e:
            print(f"❌ 导出 PDF 失败: {e}")
            self.toast_message_signal.emit("❌ 导出 PDF 失败，请检查 wkhtmltopdf 是否安装！")

    @pyqtSlot(str)
    def toggle_pin_note(self, note_id):
        try:
            order_data = read_json_file(ORDER_FILE, {"pins": [], "order": []})
            if note_id in order_data["pins"]:
                order_data["pins"].remove(note_id)
            else:
                order_data["pins"].insert(0, note_id)
            write_json_file(ORDER_FILE, order_data)
            self.get_all_notes()
            self.toast_message_signal.emit("📌 笔记已置顶！" if note_id in order_data["pins"] else "⛔ 笔记已取消置顶！")
        except Exception as e:
            print(f"置顶笔记失败: {e}")

    @pyqtSlot()
    def export_backup(self):
        try:
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            zip_name = f"backup_岸灯鸢花_{timestamp}.zip"
            zip_path = get_abs_path(zip_name)
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for folder in [NOTES_DIR, IMAGES_DIR, ATTACHMENTS_DIR]:
                    for root, dirs, files in os.walk(get_abs_path(folder)):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arcname = os.path.relpath(file_path, os.getcwd())
                            zipf.write(file_path, arcname)
            self.toast_message_signal.emit(f"备份已导出: {zip_name}")
        except Exception as e:
            print(f"导出备份失败: {e}")
            self.toast_message_signal.emit("导出备份失败！")

    @pyqtSlot()
    def get_all_notes(self):
        try:
            files = glob.glob(get_abs_path(NOTES_DIR, "*.md"))
            notes = []
            pinned_ids = []
            order_map = {}
            order_data = read_json_file(ORDER_FILE, {"pins": [], "order": []})
            if isinstance(order_data, list):
                order_map = {oid: idx for idx, oid in enumerate(order_data)}
            else:
                pinned_ids = order_data.get("pins", [])
                order_map = {oid: idx for idx, oid in enumerate(order_data.get("order", []))}

            for fp in files:
                note_id = os.path.splitext(os.path.basename(fp))[0]
                with open(fp, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                    title = "无标题笔记"
                    if lines and lines[0].startswith('# '):
                        title = lines[0][2:].strip()
                mtime = os.path.getmtime(fp)
                date_prefix = datetime.datetime.fromtimestamp(mtime).strftime("%Y-%m-%d")
                note_data = {"id": note_id, "title": title, "date": date_prefix}
                if note_id in pinned_ids:
                    note_data["pinned"] = True
                notes.append(note_data)

            # 过滤分离中的笔记
            if self.detached_note_ids:
                notes = [n for n in notes if n['id'] not in self.detached_note_ids]

            if not notes and not self.detached_note_ids:
                default_id = str(int(time.time()))
                with open(get_abs_path(NOTES_DIR, f"{default_id}.md"), 'w', encoding='utf-8') as f:
                    f.write("# 欢迎使用\n\n> 按设置的快捷键截图试试吧！")
                notes.append({"id": default_id, "title": "欢迎使用", "date": datetime.datetime.now().strftime("%Y-%m-%d")})

            notes.sort(key=lambda x: order_map.get(x['id'], float('inf')))
            self.update_sidebar_signal.emit(json.dumps(notes))
        except Exception as e:
            print(f"读取笔记目录出错: {e}")

    @pyqtSlot(str, str)
    def reorder_notes(self, dragged_id, target_id):
        try:
            order_data = read_json_file(ORDER_FILE, {"pins": [], "order": []})
            if dragged_id in order_data["pins"]:
                target_list = order_data["pins"]
            else:
                target_list = order_data["order"]
            if dragged_id in target_list and target_id in target_list:
                target_list.remove(dragged_id)
                target_idx = target_list.index(target_id)
                target_list.insert(target_idx, dragged_id)
                write_json_file(ORDER_FILE, order_data)
                self.get_all_notes()
        except Exception as e:
            print(f"排序保存失败: {e}")

    @pyqtSlot(str, str, str)
    def save_note(self, note_id, title, content):
        try:
            stripped = content.strip()
            if not stripped or not re.match(r'^#+\s', stripped):
                content = f"# {title}\n\n" + content
            file_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            self._clean_unused_images()
            self._sidebar_refresh_timer.start(300)
        except Exception as e:
            print(f"保存笔记出错: {e}")
            self.toast_message_signal.emit("保存笔记失败！")

    def _clean_unused_images(self):
        try:
            all_used = set()
            all_md = glob.glob(get_abs_path(NOTES_DIR, "*.md"))
            img_pattern = r'!\[.*?\]\((.*?)\)'
            for mf in all_md:
                with open(mf, 'r', encoding='utf-8') as f:
                    md_content = f.read()
                    for p in re.findall(img_pattern, md_content):
                        base_name = os.path.basename(p)
                        if 'images' in p or base_name.startswith(('screenshot_', 'dropped_', 'pasted_')):
                            all_used.add(base_name)
            images_dir = get_abs_path(IMAGES_DIR)
            if os.path.exists(images_dir):
                for f in os.listdir(images_dir):
                    if f.endswith(('.png', '.jpg')) and f not in all_used:
                        safe_remove(os.path.join(images_dir, f))
        except Exception as e:
            print(f"清理截图时发生意外: {e}")

    @pyqtSlot(str)
    def get_note_content(self, note_id):
        if not note_id: return
        try:
            file_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                lines = content.split('\n')
                title = "无标题笔记"
                if lines and lines[0].startswith('# '):
                    title = lines[0][2:].strip()
                self.load_note_signal.emit(title, content)
            else:
                self.load_note_signal.emit("笔记丢失", "该笔记文件已被移除。")
        except Exception as e:
            print(f"获取笔记内容出错: {e}")

    @pyqtSlot(str)
    def duplicate_note(self, note_id):
        try:
            source_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            if os.path.exists(source_path):
                with open(source_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                lines = content.split('\n')
                title = "无标题笔记"
                if lines and lines[0].startswith('# '):
                    title = lines[0][2:].strip() + " (副本)"
                new_id = str(int(time.time() * 1000)) + str(random.randint(100, 999))
                new_path = get_abs_path(NOTES_DIR, f"{new_id}.md")
                if lines and lines[0].startswith('# '):
                    lines[0] = f"# {title}\n"
                    content = "\n".join(lines)
                else:
                    content = f"# {title}\n\n" + content
                with open(new_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                self.get_all_notes()
                return new_id
        except Exception as e:
            print(f"复制笔记出错: {e}")
            return None

    @pyqtSlot(str)
    def delete_note(self, note_id):
        try:
            file_path = get_abs_path(NOTES_DIR, f"{note_id}.md")
            if not os.path.exists(file_path): return
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            note_trash_folder = get_abs_path(TRASH_DIR, note_id)
            os.makedirs(note_trash_folder, exist_ok=True)
            shutil.move(file_path, os.path.join(note_trash_folder, f"{note_id}.md"))

            img_pattern = r'!\[.*?\]\((.*?)\)'
            att_pattern = r'\[📎 .*?\]\((.*?)\)'
            current_files = set()
            for p in re.findall(img_pattern, content): current_files.add(os.path.basename(p))
            for p in re.findall(att_pattern, content): current_files.add(os.path.basename(p))

            other_mds = glob.glob(get_abs_path(NOTES_DIR, "*.md"))
            used_elsewhere = set()
            for mf in other_mds:
                if mf == file_path: continue
                with open(mf, 'r', encoding='utf-8') as f:
                    other_content = f.read()
                    for p in re.findall(img_pattern, other_content): used_elsewhere.add(os.path.basename(p))
                    for p in re.findall(att_pattern, other_content): used_elsewhere.add(os.path.basename(p))

            for f_name in current_files:
                full_img = get_abs_path(IMAGES_DIR, f_name)
                full_att = get_abs_path(ATTACHMENTS_DIR, f_name)
                staging_path = get_abs_path(STAGING_DIR, f_name)
                if f_name not in used_elsewhere:
                    safe_remove(full_img)
                    safe_remove(full_att)
                    safe_remove(staging_path)
                else:
                    if os.path.exists(full_img):
                        shutil.move(full_img, os.path.join(note_trash_folder, f_name))
                    if os.path.exists(full_att):
                        shutil.move(full_att, os.path.join(note_trash_folder, f_name))
                    if os.path.exists(staging_path):
                        shutil.move(staging_path, os.path.join(note_trash_folder, f_name))

            trash = read_json_file(TRASH_FILE, [])
            trash.append({"id": note_id, "title": "已删除笔记", "content": content, "deleted_at": int(time.time())})
            while sys.getsizeof(json.dumps(trash)) > MAX_TRASH_SIZE:
                trash.pop(0)
            write_json_file(TRASH_FILE, trash)

            # 清理分离状态
            self.detached_note_ids.discard(note_id)

            self.get_all_notes()
        except Exception as e:
            print(f"删除笔记出错: {e}")

    @pyqtSlot()
    def get_trash_items(self):
        if not os.path.exists(TRASH_FILE):
            self.update_trash_signal.emit("[]")
            return
        trash = read_json_file(TRASH_FILE, [])
        self.update_trash_signal.emit(json.dumps(trash))

    @pyqtSlot(str)
    def restore_note(self, note_id):
        try:
            trash = read_json_file(TRASH_FILE, [])
            restored = None
            new_trash = []
            for item in trash:
                if item['id'] == note_id:
                    restored = item
                else:
                    new_trash.append(item)
            if restored:
                note_trash_folder = get_abs_path(TRASH_DIR, note_id)
                if os.path.exists(note_trash_folder):
                    shutil.move(os.path.join(note_trash_folder, f"{note_id}.md"),
                                get_abs_path(NOTES_DIR, f"{note_id}.md"))
                    for f in os.listdir(note_trash_folder):
                        if f.endswith(('.png', '.jpg')):
                            shutil.move(os.path.join(note_trash_folder, f), get_abs_path(IMAGES_DIR, f))
                        elif not f.endswith('.md'):
                            shutil.move(os.path.join(note_trash_folder, f), get_abs_path(ATTACHMENTS_DIR, f))
                    os.rmdir(note_trash_folder)
                write_json_file(TRASH_FILE, new_trash)
                self.get_trash_items()
                self.get_all_notes()
        except Exception as e:
            print(f"还原笔记出错: {e}")

    @pyqtSlot(str)
    def clear_trash_permanently(self, note_id):
        try:
            note_trash_folder = get_abs_path(TRASH_DIR, note_id)
            if os.path.exists(note_trash_folder):
                shutil.rmtree(note_trash_folder)
            trash = read_json_file(TRASH_FILE, [])
            new_trash = [item for item in trash if item['id'] != note_id]
            write_json_file(TRASH_FILE, new_trash)
            self.get_trash_items()
            self.toast_message_signal.emit("已彻底从硬盘删除该笔记！")
        except Exception as e:
            print(f"永久删除出错: {e}")

    @pyqtSlot()
    def clear_trash(self):
        try:
            safe_remove(TRASH_FILE)
            if os.path.exists(TRASH_DIR):
                shutil.rmtree(TRASH_DIR)
                os.makedirs(TRASH_DIR)
            self.get_trash_items()
        except Exception as e:
            print(f"清空回收站出错: {e}")

    @pyqtSlot()
    def get_system_fonts(self):
        try:
            if hasattr(self, '_cached_fonts'):
                self.load_fonts_signal.emit(json.dumps(self._cached_fonts, ensure_ascii=False))
                return

            font_set = set()
            if is_windows():
                font_set = _get_windows_font_families()  # ← 替换原来整段注册表代码
            else:
                from PyQt6.QtGui import QFontDatabase
                font_set = set(QFontDatabase().families())

            font_list = sorted([f for f in font_set if f and f.strip()])
            self._cached_fonts = font_list
            self.load_fonts_signal.emit(json.dumps(font_list, ensure_ascii=False))
        except Exception:
            fallback = ["Microsoft YaHei", "SimHei", "Arial", "Times New Roman", "Segoe UI"]
            self.load_fonts_signal.emit(json.dumps(fallback, ensure_ascii=False))

    @pyqtSlot()
    def get_theme_packs(self):
        try:
            theme_dir = get_abs_path("themes")
            packs = []
            if os.path.exists(theme_dir):
                for f in glob.glob(os.path.join(theme_dir, "*.json")):
                    try:
                        with open(f, 'r', encoding='utf-8') as fp:
                            pack = json.load(fp)
                            if "name" in pack and "vars" in pack:
                                packs.append(pack)
                    except Exception as e:
                        print(f"❌ 读取主题包失败 {f}: {e}")
            self.theme_packs_signal.emit(json.dumps(packs))
        except Exception as e:
            print(f"❌ 主题包整体加载失败: {e}")
            self.theme_packs_signal.emit("[]")

    @pyqtSlot(str)
    def update_hotkey(self, new_key):
        # 标准化输入，去空格、转小写
        clean = new_key.strip().lower().replace(' ', '')
        seq = QKeySequence(clean)
        if seq.isEmpty():
            seq = QKeySequence("ctrl+shift+s")
        # 删除旧快捷键
        if hasattr(self.parent_window, 'hotkey_shortcut'):
            try:
                self.parent_window.hotkey_shortcut.setEnabled(False)
                self.parent_window.hotkey_shortcut.deleteLater()
            except Exception:
                pass
        # 创建新快捷键
        self.parent_window.hotkey_shortcut = QShortcut(seq, self.parent_window)
        self.parent_window.hotkey_shortcut.activated.connect(self.parent_window.screenshot_signal.emit)
        self.parent_window.hotkey_shortcut.setContext(Qt.ShortcutContext.ApplicationShortcut)
        # 反馈
        self.load_fonts_signal.emit("OK")

    @pyqtSlot()
    def toggle_floating_top(self):
        self.parent_window.toggle_floating_top()

    @pyqtSlot()
    def toggle_always_on_top(self):
        self.parent_window.toggle_always_on_top()

    @pyqtSlot(str)
    def open_floating_window(self, note_id):
        QTimer.singleShot(0, lambda: self._open_float(note_id))

    def _open_float(self, note_id):
        for fw in self.parent_window.floating_windows:
            if getattr(fw, 'note_id', None) == note_id:
                fw.raise_()
                fw.activateWindow()
                return
        self.detached_note_ids.add(note_id)
        self.parent_window.create_floating_window_with_id(note_id)
        self.get_all_notes()

    def attach_note(self, note_id=None):
        if note_id:
            self.detached_note_ids.discard(note_id)
        else:
            self.detached_note_ids.clear()
        self.get_all_notes()

    @pyqtSlot(str)
    def receive_image(self, path):
        self.image_saved_signal.emit(path.replace("\\", "/"))

# ---------- 浮窗窗口 ----------
class FloatingWindow(QMainWindow):
    def __init__(self, backend, channel, initial_pos, note_id, main_window):
        super().__init__()
        self.main_window = main_window
        self.note_id = note_id
        self.setWindowTitle("浮窗笔记")
        self.resize(700, 550)
        self.setWindowFlags(Qt.WindowType.Window)
        if is_windows():
            try:
                hwnd = int(self.winId())
                DwmSetWindowAttribute = ctypes.windll.dwmapi.DwmSetWindowAttribute
                white_color = ctypes.c_int(0x00FFFFFF)
                DwmSetWindowAttribute(hwnd, 0x35, ctypes.byref(white_color), ctypes.sizeof(white_color))
            except Exception:
                pass
        self.move(initial_pos.x() - 100, initial_pos.y() - 50)
        self.browser = QWebEngineView()
        self.browser.setPage(CustomWebEnginePage(QWebEngineProfile.defaultProfile(), self.browser))
        self.browser.setStyleSheet("background: #ffffff;")
        self.backend = backend
        self.local_channel = QWebChannel()
        self.local_channel.registerObject("backend", self.backend)
        self.float_actions = FloatActions(self)
        self.local_channel.registerObject("floatActions", self.float_actions)
        self.browser.page().setWebChannel(self.local_channel)
        float_path = get_abs_path("web", "float.html")
        self.browser.setUrl(QUrl(f"{QUrl.fromLocalFile(float_path).toString()}?id={note_id}"))
        self.setCentralWidget(self.browser)

    def toggle_float_top(self):
        flags = self.windowFlags()
        if flags & Qt.WindowType.WindowStaysOnTopHint:
            self.setWindowFlags(flags & ~Qt.WindowType.WindowStaysOnTopHint)
        else:
            self.setWindowFlags(flags | Qt.WindowType.WindowStaysOnTopHint)
        self.show()

    def open_image_editor(self, image_path):
        local_path = QUrl(image_path).toLocalFile()
        abs_path = os.path.abspath(local_path)
        if not os.path.exists(abs_path):
            return
        dialog = ImageEditorDialog(abs_path, self)
        dialog.exec()

    def closeEvent(self, event):
        self.browser.page().runJavaScript("""
            (function() {
                if (typeof currentNoteId !== 'undefined' && window.backend) {
                    const title = document.getElementById('noteTitleInput').value || '无标题笔记';
                    const content = document.getElementById('editorInput').value;
                    window.backend.save_note(currentNoteId, title, content);
                }
            })();
        """)
        if self.main_window:
            # 通知后端该笔记回归
            self.main_window.backend.attach_note(self.note_id)
            self.main_window.remove_floating_window(self)
        event.accept()

# ---------- 主窗口 ----------
class MainWindow(QMainWindow):
    screenshot_signal = pyqtSignal()

    def __init__(self):
        super().__init__()
        self.setWindowTitle("岸灯鸢花 - 便签")
        self.resize(960, 680)
        icon_path = get_abs_path("icon.png")
        if os.path.exists(icon_path):
            app_icon = QIcon(icon_path)
            self.setWindowIcon(app_icon)
            app = QApplication.instance()
            if app:
                app.setWindowIcon(app_icon)
        self.floating_windows = []
        self.init_tray()

        profile = QWebEngineProfile.defaultProfile()
        settings = profile.settings()
        settings.setAttribute(QWebEngineSettings.WebAttribute.LocalContentCanAccessFileUrls, True)
        settings.setAttribute(QWebEngineSettings.WebAttribute.LocalContentCanAccessRemoteUrls, True)
        settings.setAttribute(QWebEngineSettings.WebAttribute.JavascriptEnabled, True)

        self.browser = QWebEngineView()
        self.browser.setPage(CustomWebEnginePage(profile, self.browser))
        html_path = get_abs_path("web", "index.html")
        self.browser.setUrl(QUrl.fromLocalFile(html_path))

        central_widget = QWidget()
        layout = QVBoxLayout(central_widget)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(self.browser)
        self.setCentralWidget(central_widget)

        self.channel = QWebChannel()
        self.backend = Backend(self)
        self.channel.registerObject("backend", self.backend)
        self.browser.page().setWebChannel(self.channel)

        # 暂时不注册全局热键
        self.screenshot_signal.connect(self.take_screenshot)
        from PyQt6.QtGui import QKeySequence, QShortcut
        self.hotkey_shortcut = QShortcut(QKeySequence("Ctrl+Shift+S"), self)
        self.hotkey_shortcut.activated.connect(self.screenshot_signal.emit)
        self.hotkey_shortcut.setContext(Qt.ShortcutContext.ApplicationShortcut)

    def init_tray(self):
        self.tray_icon = QSystemTrayIcon(self)
        icon_path = get_abs_path("icon.png")
        if os.path.exists(icon_path):
            self.tray_icon.setIcon(QIcon(icon_path))
        tray_menu = QMenu()
        show_action = QAction("显示主窗口", self)
        quit_action = QAction("彻底退出", self)
        show_action.triggered.connect(self.show_normal)
        quit_action.triggered.connect(self.quit_app)
        tray_menu.addAction(show_action)
        tray_menu.addAction(quit_action)
        self.tray_icon.setContextMenu(tray_menu)
        self.tray_icon.activated.connect(self.tray_icon_activated)
        self.tray_icon.show()

    def tray_icon_activated(self, reason):
        if reason == QSystemTrayIcon.ActivationReason.DoubleClick:
            self.show_normal()

    def show_normal(self):
        self.showNormal()
        self.raise_()
        self.activateWindow()

    def quit_app(self):
        self.tray_icon.hide()
        QApplication.quit()

    def toggle_always_on_top(self):
        flags = self.windowFlags()
        if flags & Qt.WindowType.WindowStaysOnTopHint:
            self.setWindowFlags(flags ^ Qt.WindowType.WindowStaysOnTopHint)
        else:
            self.setWindowFlags(flags | Qt.WindowType.WindowStaysOnTopHint)
        self.show()

    def remove_floating_window(self, window):
        QTimer.singleShot(100, lambda: self._safe_remove(window))

    def _safe_remove(self, window):
        if window in self.floating_windows:
            self.floating_windows.remove(window)

    def create_floating_window_with_id(self, note_id):
        if not note_id: return
        cursor_pos = QCursor.pos()
        window = FloatingWindow(self.backend, self.channel, cursor_pos, note_id, self)
        window.show()
        self.floating_windows.append(window)

    def closeEvent(self, event):
        event.ignore()
        self.hide()
        self.tray_icon.showMessage("岸灯鸢花", "软件已最小化到系统托盘，双击可恢复。",
                                   QSystemTrayIcon.MessageIcon.Information, 2000)

    @pyqtSlot()
    def take_screenshot(self):
        QTimer.singleShot(30, self._start_region_capture)

    def _start_region_capture(self):
        self.capture_widget = ScreenCaptureWidget()
        self.capture_widget.region_selected.connect(self._on_region_selected)
        self.capture_widget.show()

    def _on_region_selected(self, rect):
        screen = QApplication.screenAt(rect.center()) or QApplication.primaryScreen()
        if screen:
            pixmap = screen.grabWindow(0, rect.x(), rect.y(), rect.width(), rect.height())
            filename = f"screenshot_{int(time.time())}.png"
            full_path = get_abs_path(IMAGES_DIR, filename)
            os.makedirs(os.path.dirname(full_path), exist_ok=True)
            pixmap.save(full_path, "PNG")
            self.backend.receive_image(full_path)
        self.show()

# ---------- 主程序入口 ----------
if __name__ == "__main__":
    app = QApplication(sys.argv)

    shared_mem = QSharedMemory("AndengYuanhua_Unique_Instance")
    if shared_mem.attach():
        sys.exit(0)
    else:
        shared_mem.create(1)

    app.setQuitOnLastWindowClosed(False)
    window = MainWindow()
    screen = app.primaryScreen()
    screen_geo = screen.availableGeometry()
    x = (screen_geo.width() - window.width()) // 2
    y = (screen_geo.height() - window.height()) // 2
    window.move(x, y)
    window.show()
    sys.exit(app.exec())